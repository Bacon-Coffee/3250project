"""Phase 10: presentation slide deck.

Builds a 16:9 PPTX deck that walks through the project's narrative
(motivation -> methodology -> five figures -> conclusion) reusing the
PNGs produced by ``make_figures.py``. The deck is meant to be opened
in Keynote / LibreOffice / PowerPoint and presented as-is.

Run:
    python -m src.experiments.make_slides
Output:
    results/figures/slides.pptx
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[2]
FIG_DIR = REPO_ROOT / "results" / "figures"
OUT_PATH = FIG_DIR / "slides.pptx"

# 16:9 widescreen, in inches.
SLIDE_W = 13.333
SLIDE_H = 7.5

NAVY = RGBColor(0x0B, 0x25, 0x45)
GOLD = RGBColor(0xFF, 0xD1, 0x66)
SOFT = RGBColor(0x33, 0x33, 0x33)
ACCENT = RGBColor(0xD6, 0x28, 0x28)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


@dataclass
class FigSlide:
    title: str
    image: str
    bullets: list[str]


FIG_SLIDES: list[FigSlide] = [
    FigSlide(
        title="Fig 1.  P99 response and Jain's fairness",
        image="fig1_p99_fairness.png",
        bullets=[
            "Six schedulers x three workloads, three seeds each.",
            "CFS and EEVDF cut P99 by 2-3 orders of magnitude vs FCFS / SJF.",
            "Jain's index is uniform: aggregate fairness saturates once no "
            "task starves -- the real story is in windowed lag.",
        ],
    ),
    FigSlide(
        title="Fig 2.  CFS vs EEVDF on Zijlstra's case",
        image="fig2_cfs_vs_eevdf.png",
        bullets=[
            "Reproduces the short-task-vs-long-task scenario from the "
            "Linux 6.6 EEVDF patch series.",
            "EEVDF's lag-based wakeup keeps the short job's waiting time "
            "bounded; CFS's pure vruntime ordering does not.",
            "Mean waiting time falls ~2% on nice-mixed load; tail P99 "
            "behavior depends on slice and latency_nice.",
        ],
    ),
    FigSlide(
        title="Fig 3.  RB-tree microbench",
        image="fig3_rbtree_microbench.png",
        bullets=[
            "Hand-written CLRS RB tree vs sortedcontainers.SortedList.",
            "2-3x per-operation overhead, dominated by Python interpreter "
            "cost on rotation fixups.",
            "Justified: every insert/delete path in src/core/rbtree.py "
            "maps to a named case in lib/rbtree.c (Methodology asset).",
        ],
    ),
    FigSlide(
        title="Fig 4.  Bitbrains GWA-T-12 CDFs",
        image="fig4_bitbrains_cdf.png",
        bullets=[
            "Real production VM trace, mapped to per-process bursts "
            "(caveat: VM-level aggregation declared in Methodology).",
            "CFS / EEVDF dominate the upper tail; RR and MLFQ trade tail "
            "for steady throughput.",
            "Shape matches the Lozi 2016 'cores left idle' observation in "
            "the no-balancer regime.",
        ],
    ),
    FigSlide(
        title="Fig 5.  Dual-core idle balance",
        image="fig5_dual_core_lb.png",
        bullets=[
            "All tasks initially targeted at CPU0; CPU1 starves without a "
            "balancer (top panel).",
            "Idle-balance enabled (bottom): CPU1 climbs from 0% to 97.3% "
            "while CPU0 stays at 100%.",
            "Migration cost = 1 tick per stolen task; cache-thrash regime "
            "not modeled (SMP-coherence caveat).",
        ],
    ),
]

TITLE_SLIDE = {
    "title": "From FCFS to EEVDF",
    "subtitle": "An Empirical Reconstruction of the Linux Scheduler\n"
                "Evolution with a Hand-Written Red-Black Tree and\n"
                "Dual-Core Load Balancing",
    "footer": "Course Project, CS 3250  |  Spring 2026",
}

MOTIVATION_BULLETS = [
    "Linux replaced CFS with EEVDF in v6.6 (Oct 2023), ending a 16-year "
    "CFS era introduced by Ingo Molnar in v2.6.23 (2007).",
    "Textbooks rarely connect FCFS / SJF / RR / MLFQ to CFS, and almost "
    "never to EEVDF.",
    "We rebuild the full path in one event-driven Python simulator with "
    "a unified ABC interface and a hand-written RB tree shared by CFS "
    "and EEVDF.",
    "Hypotheses: H1 fairness, H2 RB-tree overhead, H3 EEVDF tail "
    "latency, H4 dual-core load balancing.",
]

METHOD_BULLETS = [
    "Unified scheduler ABC: on_arrival / on_tick / pick_next / "
    "peek_steal_candidate -- no algorithm-specific branches.",
    "From-scratch CLRS red-black tree (left/right rotations, double-red "
    "and double-black fixup) reused by CFS and EEVDF; only the key "
    "function differs (vruntime vs virtual_deadline).",
    "Kernel constants traced to source: NICE_0_LOAD=1024, "
    "prio_to_weight[], sched_latency, min_granularity.",
    "Dual-CPU model: each CPU owns a runqueue; idle-balance triggers "
    "steal of the rightmost candidate; +1 tick migration cost.",
    "Workloads: CPU-heavy, I/O-heavy, mixed, nice-weighted (3 seeds) + "
    "Bitbrains GWA-T-12 trace.",
]

CONCLUSION_BULLETS = [
    "H1 (fairness): aggregate Jain's saturates; report windowed lag.",
    "H2 (RB-tree overhead crossover): not observed at <=500 tasks.",
    "H3 (EEVDF tail latency): supported on nice-mixed load (~2% mean "
    "waiting-time cut vs CFS).",
    "H4 (dual-core balancing): supported -- CPU1 climbs 0% -> 97.3% "
    "with idle-balance alone.",
    "Future work: O(log n) augmented-tree eligibility scan; four-CPU "
    "two-level sched_domain; Google cluster trace at process granularity.",
]


def _add_background(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                Inches(SLIDE_W), Inches(SLIDE_H))
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False


def _add_text(slide, left, top, width, height, text, *,
              size=18, bold=False, color=SOFT, align_center=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    lines = [text] if isinstance(text, str) else list(text)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        if align_center:
            p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return box


def _add_bullets(slide, left, top, width, height, bullets,
                 *, size=18, color=SOFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"•  {bullet}"
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color


def _title_bar(slide, text: str) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                 Inches(SLIDE_W), Inches(0.9))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.shadow.inherit = False
    _add_text(slide, Inches(0.5), Inches(0.15), Inches(SLIDE_W - 1),
              Inches(0.7), text, size=28, bold=True, color=WHITE)


def _add_image_fitted(slide, image_path: Path,
                      left: float, top: float,
                      max_w: float, max_h: float) -> None:
    with Image.open(image_path) as im:
        iw, ih = im.size
    aspect = iw / ih
    if max_w / max_h > aspect:
        h = max_h
        w = h * aspect
    else:
        w = max_w
        h = w / aspect
    x = left + (max_w - w) / 2
    y = top + (max_h - h) / 2
    slide.shapes.add_picture(str(image_path), Inches(x), Inches(y),
                             Inches(w), Inches(h))


def _build_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, NAVY)
    _add_text(slide, Inches(0.5), Inches(2.4), Inches(SLIDE_W - 1),
              Inches(1.3), TITLE_SLIDE["title"],
              size=64, bold=True, color=GOLD, align_center=True)
    _add_text(slide, Inches(1.0), Inches(3.9), Inches(SLIDE_W - 2),
              Inches(2.0), TITLE_SLIDE["subtitle"].split("\n"),
              size=26, color=WHITE, align_center=True)
    _add_text(slide, Inches(1.0), Inches(6.5), Inches(SLIDE_W - 2),
              Inches(0.6), TITLE_SLIDE["footer"],
              size=18, color=GOLD, align_center=True)


def _build_text_slide(prs: Presentation, title: str,
                      bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _title_bar(slide, title)
    _add_bullets(slide, Inches(0.7), Inches(1.3),
                 Inches(SLIDE_W - 1.4), Inches(SLIDE_H - 1.5),
                 bullets, size=22)


def _build_figure_slide(prs: Presentation, fig: FigSlide) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _title_bar(slide, fig.title)
    img_path = FIG_DIR / fig.image
    if img_path.exists():
        _add_image_fitted(slide, img_path,
                          left=0.4, top=1.2,
                          max_w=7.6, max_h=5.9)
    else:
        _add_text(slide, Inches(0.4), Inches(3.5), Inches(7.6),
                  Inches(0.6), f"[missing image: {fig.image}]",
                  size=20, color=ACCENT, align_center=True)
    _add_bullets(slide, Inches(8.3), Inches(1.4),
                 Inches(SLIDE_W - 8.6), Inches(SLIDE_H - 1.6),
                 fig.bullets, size=18)


def _build_closing(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, NAVY)
    _add_text(slide, Inches(0.5), Inches(0.5), Inches(SLIDE_W - 1),
              Inches(1.0), "Thank you", size=48, bold=True,
              color=GOLD, align_center=True)
    _add_text(slide, Inches(1.0), Inches(2.2), Inches(SLIDE_W - 2),
              Inches(0.6),
              "Reproduce with:  python -m src.experiments.run_all",
              size=22, color=WHITE, align_center=True)
    _add_text(slide, Inches(1.0), Inches(3.2), Inches(SLIDE_W - 2),
              Inches(0.6),
              "Key references:  Molnar 2007 (CFS), Zijlstra 2023 (EEVDF), "
              "Lozi 2016 (CFS bugs), CLRS Ch. 13 (RB tree).",
              size=18, color=WHITE, align_center=True)


def _set_slide_size(prs: Presentation) -> None:
    prs.slide_width = Emu(int(SLIDE_W * 914400))
    prs.slide_height = Emu(int(SLIDE_H * 914400))


def build_slides() -> Path:
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    _set_slide_size(prs)

    _build_title(prs)
    _build_text_slide(prs, "Motivation", MOTIVATION_BULLETS)
    _build_text_slide(prs, "Methodology", METHOD_BULLETS)
    for fig in FIG_SLIDES:
        _build_figure_slide(prs, fig)
    _build_text_slide(prs, "Conclusion & Future Work", CONCLUSION_BULLETS)
    _build_closing(prs)

    prs.save(OUT_PATH)
    return OUT_PATH


if __name__ == "__main__":
    out = build_slides()
    print(f"slides written: {out}")
